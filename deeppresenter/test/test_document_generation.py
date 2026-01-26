"""
Test document generation with python-pptx.
"""

import pytest


class TestDocumentGeneration:
    """Test document generation with python-pptx."""

    @pytest.mark.asyncio
    async def test_python_pptx_comprehensive(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test python-pptx: multiple slides + text + shapes + chinese content."""
        script_content = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title slide with chinese
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]
title.text = "Python-PPTX 综合测试"
subtitle.text = "自动生成演示文稿\\n沙箱环境测试"

# Slide 2: Content slide with bullets
content_slide = prs.slides.add_slide(prs.slide_layouts[1])
title = content_slide.shapes.title
content = content_slide.placeholders[1]
title.text = "主要功能特性"
tf = content.text_frame
tf.text = "功能 1: 基础文本支持"
p = tf.add_paragraph()
p.text = "功能 2: 多级列表"
p.level = 1
p = tf.add_paragraph()
p.text = "功能 3: 中文字体渲染"
p.level = 1
p = tf.add_paragraph()
p.text = "功能 4: 自定义样式"

# Slide 3: Blank slide with shapes
blank_slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = blank_slide.shapes

# Add title
title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "图形与形状演示"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Add rectangle
rect = shapes.add_shape(
    1,  # Rectangle
    Inches(1), Inches(2), Inches(3), Inches(2)
)
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0, 128, 255)
rect_text = rect.text_frame
rect_text.text = "蓝色矩形"

# Add circle (oval)
circle = shapes.add_shape(
    9,  # Oval
    Inches(5), Inches(2), Inches(2), Inches(2)
)
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(255, 128, 0)
circle_text = circle.text_frame
circle_text.text = "圆形"

# Slide 4: Summary
summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
title = summary_slide.shapes.title
content = summary_slide.placeholders[1]
title.text = "总结"
tf = content.text_frame
tf.text = "✓ 成功创建多页演示文稿"
p = tf.add_paragraph()
p.text = "✓ 支持中文内容"
p = tf.add_paragraph()
p.text = "✓ 包含图形和形状"
p = tf.add_paragraph()
p.text = "✓ 自定义样式和颜色"

# Save presentation
prs.save('综合演示.pptx')
print(f'综合演示.pptx created ({len(prs.slides)} slides)')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_pptx.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_pptx.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"python-pptx test failed: {exec_result.text}"
        assert "综合演示.pptx created" in exec_result.text
        assert "4 slides" in exec_result.text

        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "综合演示.pptx" in list_result.text
